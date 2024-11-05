from pptx import Presentation

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]

title.text = "LoRa-Based AGV Tracking in Manufacturing"
subtitle.text = "Tracking AGVs without GPS or Wi-Fi using LoRa Communication"

# Slide 2: Introduction
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
title_2 = slide_2.shapes.title
title_2.text = "Introduction"

content_2 = slide_2.shapes.placeholders[1].text_frame
content_2.text = "In a manufacturing environment, tracking AGVs (Automated Guided Vehicles) without relying on GPS or Wi-Fi is possible using LoRa technology."

p = content_2.add_paragraph()
p.text = "LoRa provides long-range, low-power communication ideal for industrial settings."

# Slide 3: LoRa System Overview
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
title_3 = slide_3.shapes.title
title_3.text = "LoRa System Overview"

content_3 = slide_3.shapes.placeholders[1].text_frame
content_3.text = "A basic LoRa system for AGV tracking includes:"

p1 = content_3.add_paragraph()
p1.text = "1. LoRa transmitter on the AGV (Node)"

p2 = content_3.add_paragraph()
p2.text = "2. LoRa receiver (Gateway) at a centralized location"

p3 = content_3.add_paragraph()
p3.text = "3. Optional LoRa antenna for better range and coverage"

# Slide 4: LoRa Node on AGV
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
title_4 = slide_4.shapes.title
title_4.text = "LoRa Node on AGV"

content_4 = slide_4.shapes.placeholders[1].text_frame
content_4.text = "The AGV's PLC communicates with the LoRa transmitter to send real-time position and status data to the gateway."

p4 = content_4.add_paragraph()
p4.text = "Modules like SX1276 or RFM95 are typically used for LoRa transmission."

# Slide 5: LoRa Gateway
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
title_5 = slide_5.shapes.title
title_5.text = "LoRa Gateway at Central Location"

content_5 = slide_5.shapes.placeholders[1].text_frame
content_5.text = "The LoRa gateway receives the AGV data and forwards it to the central control system or HMI."

p5 = content_5.add_paragraph()
p5.text = "Gateways like RAK7243C or Dragino LG01 can be used."

# Slide 6: Advantages of LoRa for AGV Tracking
slide_6 = prs.slides.add_slide(prs.slide_layouts[1])
title_6 = slide_6.shapes.title
title_6.text = "Advantages of LoRa for AGV Tracking"

content_6 = slide_6.shapes.placeholders[1].text_frame
content_6.text = "1. Long range (up to 10-15 km in open environments)"

p6 = content_6.add_paragraph()
p6.text = "2. Low power consumption, ideal for battery-powered AGVs"

p7 = content_6.add_paragraph()
p7.text = "3. Operates without GPS or Wi-Fi, reducing complexity in indoor environments"

# Slide 7: Conclusion
slide_7 = prs.slides.add_slide(prs.slide_layouts[1])
title_7 = slide_7.shapes.title
title_7.text = "Conclusion"

content_7 = slide_7.shapes.placeholders[1].text_frame
content_7.text = "LoRa is an effective solution for tracking AGVs in manufacturing environments, providing long-range, reliable communication without the need for GPS or Wi-Fi."

# Save the presentation
pptx_file = "C:/Users/pravin/OneDrive/Desktop/LoRa_AGV_Tracking.pptx"
prs.save(pptx_file)

pptx_file
